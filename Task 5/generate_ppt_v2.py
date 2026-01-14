"""Improved PPT generator with image support and simple design.

Usage:
    pip install python-pptx
    python3 generate_ppt_v2.py

Overwrites: APP_Presentation.pptx in the same folder.
"""
import re
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

MD = Path(__file__).with_name("APP_Presentation.md")
OUT = Path(__file__).with_name("APP_Presentation.pptx")
ASSETS = Path(__file__).with_name("assets")

if not MD.exists():
    raise SystemExit(f"Missing markdown file: {MD}")

text = MD.read_text(encoding="utf-8")
# Split into top-level sections (## )
sections = re.split(r"^##\s+", text, flags=re.MULTILINE)
prs = Presentation()
prs.slide_height = Inches(6)
prs.slide_width = Inches(10)

# Title slide (use a simple colored background)
if sections and sections[0].strip():
    title_text = sections[0].strip().splitlines()[0]
else:
    title_text = "Car Service Booking"
sub_text = "Presentation"
try:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    # background rectangle
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE.RECTANGLE (1)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(18, 97, 160)  # blue
    shape.line.fill.background()
    # title
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(1.5)).text_frame
    p = tx.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    # subtitle
    st = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(0.8)).text_frame
    st.paragraphs[0].text = sub_text
    st.paragraphs[0].font.size = Pt(18)
    st.paragraphs[0].font.color.rgb = RGBColor(230, 230, 230)
except Exception:
    pass

for sec in sections[1:]:
    lines = sec.strip().splitlines()
    if not lines:
        continue
    title = lines[0].strip()
    body_lines = lines[1:]
    # create a clean slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # set a light background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)

    # title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.8))
    tframe = title_box.text_frame
    tframe.paragraphs[0].text = title
    tframe.paragraphs[0].font.size = Pt(28)
    tframe.paragraphs[0].font.bold = True
    tframe.paragraphs[0].font.color.rgb = RGBColor(18, 97, 160)

    # content area
    left = Inches(0.6)
    top = Inches(1.3)
    width = Inches(5.8)
    height = Inches(4.2)

    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.clear()

    # process lines: images or bullets
    for ln in body_lines:
        ln = ln.rstrip()
        if not ln:
            continue
        # image markdown: ![alt](path)
        img_m = re.match(r"!\[.*\]\((.*)\)", ln)
        if img_m:
            img_path = img_m.group(1)
            ip = Path(img_path)
            if ip.is_absolute():
                img_file = ip
            else:
                # support markdown paths like 'assets/code.png' or just 'code.png'
                parts = ip.parts
                if parts and parts[0] == 'assets':
                    img_file = ASSETS.joinpath(*parts[1:])
                else:
                    img_file = ASSETS / ip
            if img_file.exists():
                # place picture on the right side
                pic_left = Inches(7.0)
                pic_top = Inches(1.3)
                pic_w = Inches(2.8)
                slide.shapes.add_picture(str(img_file), pic_left, pic_top, width=pic_w)
            else:
                print(f"Warning: image not found: {img_file}")
            continue
        # headings inside section
        if ln.startswith('###'):
            sub = ln.lstrip('#').strip()
            p = tf.add_paragraph()
            p.text = sub
            p.font.bold = True
            p.level = 0
            continue
        m = re.match(r"^\s*([-*]|\d+\.)\s+(.*)", ln)
        if m:
            txt = m.group(2).strip()
            p = tf.add_paragraph()
            p.text = txt
            p.level = 0
        else:
            # regular paragraph - if long, split into shorter paragraphs by '. '
            parts = [s.strip() for s in re.split(r"(?<=\.)\s+", ln) if s.strip()]
            for part in parts:
                p = tf.add_paragraph()
                p.text = part
                p.level = 0

    # try to ensure text fits: reduce font size if there are many paragraphs
    n_paras = sum(1 for _ in tf.paragraphs)
    base_size = 16
    if n_paras > 8:
        size = Pt(11)
    elif n_paras > 6:
        size = Pt(12)
    elif n_paras > 4:
        size = Pt(13)
    else:
        size = Pt(base_size)
    for p in tf.paragraphs:
        p.font.size = size
        p.font.color.rgb = RGBColor(40, 40, 40)

    # footer
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(8.8), Inches(0.3)).text_frame
    footer.paragraphs[0].text = f"Car Service Booking — {datetime.utcnow().date()}"
    footer.paragraphs[0].font.size = Pt(10)
    footer.paragraphs[0].font.italic = True

prs.save(OUT)
print(f"Created {OUT}")