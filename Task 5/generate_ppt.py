"""Generate a PPTX from the APP_Presentation.md file.

Usage:
    pip install python-pptx
    python3 generate_ppt.py

Creates: APP_Presentation.pptx in the same folder.
"""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

MD = Path(__file__).with_name("APP_Presentation.md")
OUT = Path(__file__).with_name("APP_Presentation.pptx")

if not MD.exists():
    raise SystemExit(f"Missing markdown file: {MD}")

text = MD.read_text(encoding="utf-8")
# Split into top-level sections (## )
sections = re.split(r"^##\s+", text, flags=re.MULTILINE)
# first chunk is header before first '##' (likely empty)
prs = Presentation()
# Title slide
if sections and sections[0].strip():
    title_text = sections[0].strip().splitlines()[0]
else:
    title_text = "Car Service Booking"
sub_text = "Presentation"
try:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = sub_text
except Exception:
    pass

for sec in sections[1:]:
    lines = sec.strip().splitlines()
    if not lines:
        continue
    title = lines[0].strip()
    body_lines = lines[1:]
    slide_layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    # process lines: treat lines starting with '- ' or numbered lists as bullets
    for ln in body_lines:
        ln = ln.rstrip()
        if not ln:
            continue
        # headings '###' inside section become bold lines
        if ln.startswith('###'):
            sub = ln.lstrip('#').strip()
            p = tf.add_paragraph()
            p.text = sub
            p.font.bold = True
            p.level = 0
            continue
        m = re.match(r"^\s*([-\*]|\d+\.)\s+(.*)", ln)
        if m:
            txt = m.group(2).strip()
            p = tf.add_paragraph()
            p.text = txt
            p.level = 0
        else:
            # regular paragraph - add as a normal bullet (if too long, it will wrap)
            p = tf.add_paragraph()
            p.text = ln.strip()
            p.level = 0

prs.save(OUT)
print(f"Created {OUT}")
